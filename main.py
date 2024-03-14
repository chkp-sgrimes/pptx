import pptx.shapes.autoshape
import structlog
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.exc import PackageNotFoundError
from pathlib import Path
import structlog

log: structlog = structlog.getLogger()

structlog.configure(
    processors=[
        structlog.contextvars.merge_contextvars,
        structlog.processors.add_log_level,
        structlog.processors.StackInfoRenderer(),
        structlog.dev.set_exc_info,
        structlog.processors.TimeStamper(fmt="%Y-%m-%d %H:%M:%S", utc=False),
        # structlog.dev.ConsoleRenderer(),
        structlog.processors.JSONRenderer()
    ],
    # wrapper_class=structlog.make_filtering_bound_logger(logging.NOTSET),
    context_class=dict,
    logger_factory=structlog.PrintLoggerFactory(),
    cache_logger_on_first_use=False
)


def get_paragraph_text(para: pptx.shapes.autoshape.Shape) -> str:
    text_disallowed = ["©", "[Internal Use]", "Thank you"]
    char_disallowed = ["\x0b", "\n"]
    text: str = ""

    text = ''.join([' ' if i in char_disallowed else i for i in para.text]).strip(" ")
    if any(x in text for x in text_disallowed) or not len(text):
        return None
    else:
        text = f"{text}\n"
        return text


def process_text_frames(shape_text_frames: list) -> str:
    slide_text: str = ''
    additional_text = ''

    for para in shape_text_frames:
        additional_text = get_paragraph_text(para)
        if additional_text is not None:
            slide_text = ' '.join([additional_text, slide_text])
        else:
            continue

    return slide_text


def process_groups(shape_group: list) -> str:
    slide_text: str = ''

    # subgroups may have text frames
    subgroups = [element for shape in shape_group for element in shape.shapes if
                 element.shape_type == MSO_SHAPE_TYPE.GROUP]

    # there may be text frames at the slide level
    textsubs = [element for shape in shape_group for element in shape.shapes if element.has_text_frame]

    slide_text = process_text_frames(textsubs)

    return slide_text


def process_shapes(shape_list: list) -> str:
    slide_text: str = ''

    text_frames = [_ for _ in shape_list if _.has_text_frame]
    slide_text = process_text_frames(text_frames)

    groups = [_ for _ in shape_list if _.shape_type == MSO_SHAPE_TYPE.GROUP]
    slide_text = ''.join([process_groups(groups), slide_text])

    return slide_text


def main() -> None:
    # logs: StructLogger = StructLogger()
    slide_text: str = ''
    slide_number: int = 0
    slide_name: str = ''
    files: list = [Path]
    preso: Path

    files = [f for f in Path().glob("inputs/*.pptx")]

    for preso in files:
        try:
            prs: Presentation() = Presentation('inputs/' + preso.name)
            log.info("pptx_open", file=preso.name)
        except PackageNotFoundError:
            print(f"PackageNotFoundError possible encrypted file {preso.name}")
            preso.rename(f"encrypted/{preso.name}")
            log.error("pptx_encrypted", name=preso.name)
            continue

        for slide in prs.slides:
            slide_number += 1
            log.info("slide", number=slide_number, slide_id=slide.slide_id)
            shape_list = [s for s in slide.shapes if s.has_text_frame or s.shape_type == MSO_SHAPE_TYPE.GROUP]
            slide_text = process_shapes(shape_list)
            log.info("slide_text", current_text=slide_text)

            slide_name = f"outputs/{preso.stem}_{str(slide_number).zfill(3)}.txt"
            try:
                with open(slide_name, "w") as output_file:
                    reply = output_file.write(slide_text)
                    log.info("text_write", name=slide_name)
            except FileExistsError as e:
                log.error("file_error", error=e)

            slide_text = ''


if __name__ == '__main__':
    main()
